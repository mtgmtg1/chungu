#!/usr/bin/env python3
# [Flow: Step 1 (마크다운 블록 파싱) -> Step 2 (docx/pptx/xlsx 작성) -> Step 3 (파일 경로 반환)]
import html as html_module
import re
from pathlib import Path

from docx import Document
from docx.shared import Inches, Pt, RGBColor
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as PptxInches


# ---------------------------------------------------------------------------
# Markdown block parser
# ---------------------------------------------------------------------------

def _parse_markdown_blocks(markdown: str) -> list[dict]:
    """마크다운 텍스트를 블록 단위로 파싱하여 모든 콘텐츠를 유실 없이 반환한다.

    반환 블록 형식:
      - {"type": "heading", "level": int, "text": str}
      - {"type": "paragraph", "text": str}
      - {"type": "list", "ordered": bool, "items": list[str]}
      - {"type": "table", "headers": list[str], "rows": list[list[str]]}
      - {"type": "code", "language": str, "text": str}
    """
    lines = markdown.splitlines()
    blocks: list[dict] = []
    i = 0
    while i < len(lines):
        line = lines[i]
        stripped = line.strip()

        if not stripped:
            i += 1
            continue

        # 코드 블록
        if stripped.startswith("```"):
            language = stripped[3:].strip()
            i += 1
            code_lines: list[str] = []
            while i < len(lines) and not lines[i].strip().startswith("```"):
                code_lines.append(lines[i])
                i += 1
            i += 1  # 닫는 ``` 건너뜀
            blocks.append({"type": "code", "language": language, "text": "\n".join(code_lines)})
            continue

        # 제목
        if stripped.startswith("#"):
            level = 0
            while level < len(stripped) and stripped[level] == "#":
                level += 1
            text = stripped[level:].strip()
            blocks.append({"type": "heading", "level": level, "text": text})
            i += 1
            continue

        # 표
        if stripped.startswith("|"):
            table_lines: list[str] = []
            while i < len(lines) and lines[i].strip().startswith("|"):
                table_lines.append(lines[i].strip())
                i += 1
            if len(table_lines) >= 2:
                headers = [c.strip() for c in table_lines[0].split("|")[1:-1]]
                rows = []
                for row_line in table_lines[2:]:
                    cells = [c.strip() for c in row_line.split("|")[1:-1]]
                    rows.append(cells)
                blocks.append({"type": "table", "headers": headers, "rows": rows})
            continue

        # 목록
        list_match = re.match(r"^[-*+]\s+(.*)", stripped)
        ordered_match = re.match(r"^(\d+)\.\s+(.*)", stripped)
        if list_match or ordered_match:
            items: list[str] = []
            ordered = bool(ordered_match)
            while i < len(lines):
                item_line = lines[i].strip()
                if not item_line:
                    break
                unordered_m = re.match(r"^[-*+]\s+(.*)", item_line)
                ordered_m = re.match(r"^\d+\.\s+(.*)", item_line)
                if unordered_m:
                    ordered = False
                    items.append(unordered_m.group(1))
                    i += 1
                elif ordered_m:
                    items.append(ordered_m.group(1))
                    i += 1
                else:
                    break
            blocks.append({"type": "list", "ordered": ordered, "items": items})
            continue

        # 문단
        para_lines = [line]
        i += 1
        while i < len(lines):
            next_line = lines[i]
            next_stripped = next_line.strip()
            if not next_stripped:
                break
            if (
                next_stripped.startswith("#")
                or next_stripped.startswith("```")
                or next_stripped.startswith("|")
                or re.match(r"^[-*+]\s", next_stripped)
                or re.match(r"^\d+\.\s", next_stripped)
            ):
                break
            para_lines.append(next_line)
            i += 1
        blocks.append({"type": "paragraph", "text": " ".join(para_lines).strip()})

    return blocks


def _split_text_inline(text: str) -> list[tuple[str, dict]]:
    """마크다운 인라인 서식(**굵게**, *기울임*, ~~취소선~~)을 분리한다."""
    pattern = r"(\*\*\*[^*]+\*\*\*|\*\*[^*]+\*\*|\*[^*]+\*|~~[^~]+~~)"
    parts = re.split(pattern, text)
    result = []
    for part in parts:
        if not part:
            continue
        fmt = {}
        if part.startswith("***") and part.endswith("***"):
            fmt = {"bold": True, "italic": True}
            part = part[3:-3]
        elif part.startswith("**") and part.endswith("**"):
            fmt = {"bold": True}
            part = part[2:-2]
        elif part.startswith("*") and part.endswith("*"):
            fmt = {"italic": True}
            part = part[1:-1]
        elif part.startswith("~~") and part.endswith("~~"):
            fmt = {"strike": True}
            part = part[2:-2]
        result.append((part, fmt))
    return result


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _add_docx_runs(paragraph, text: str) -> None:
    """문단에 인라인 서식을 적용한 run을 추가한다."""
    for part, fmt in _split_text_inline(text):
        run = paragraph.add_run(part)
        if fmt.get("bold"):
            run.bold = True
        if fmt.get("italic"):
            run.italic = True
        if fmt.get("strike"):
            run.font.strike = True


def markdown_to_docx(markdown: str, out_path: Path) -> Path:
    """마크다운 전체 콘텐츠를 Word 문서로 변환한다."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    blocks = _parse_markdown_blocks(markdown)
    doc = Document()

    if not blocks:
        doc.add_paragraph("변환할 콘텐츠가 없습니다.")
        doc.save(out_path)
        return out_path

    table_idx = 0
    for block in blocks:
        if block["type"] == "heading":
            doc.add_heading(block["text"], level=min(block["level"], 6))
        elif block["type"] == "paragraph":
            p = doc.add_paragraph()
            _add_docx_runs(p, block["text"])
        elif block["type"] == "list":
            style = "List Number" if block["ordered"] else "List Bullet"
            for item in block["items"]:
                p = doc.add_paragraph(style=style)
                _add_docx_runs(p, item)
        elif block["type"] == "table":
            table_idx += 1
            headers = block["headers"]
            rows = block["rows"]
            if not headers and not rows:
                continue
            cols = max(len(headers), max((len(r) for r in rows), default=1), 1)
            doc_table = doc.add_table(rows=1 + len(rows), cols=cols)
            doc_table.style = "Table Grid"
            for i, h in enumerate(headers):
                cell = doc_table.rows[0].cells[i]
                cell.text = h
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
            for r_idx, row in enumerate(rows):
                for c_idx, val in enumerate(row):
                    if c_idx < cols:
                        doc_table.rows[r_idx + 1].cells[c_idx].text = val
            doc.add_paragraph()
        elif block["type"] == "code":
            p = doc.add_paragraph()
            run = p.add_run(block["text"])
            run.font.name = "Courier New"
            run.font.size = Pt(10)
            p.paragraph_format.left_indent = Inches(0.25)

    doc.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _blocks_to_slide_text(blocks: list[dict]) -> str:
    """블록 목록을 슬라이드 본문 텍스트로 변환한다."""
    lines = []
    for block in blocks:
        if block["type"] == "paragraph":
            lines.append(block["text"])
        elif block["type"] == "heading":
            lines.append("" + ("#" * block["level"]) + " " + block["text"])
        elif block["type"] == "list":
            for idx, item in enumerate(block["items"], start=1):
                prefix = f"{idx}. " if block["ordered"] else "• "
                lines.append(prefix + item)
        elif block["type"] == "table":
            headers = " | ".join(block["headers"])
            rows = [" | ".join(r) for r in block["rows"]]
            lines.append(headers)
            lines.append("-" * len(headers))
            lines.extend(rows)
            lines.append("")
        elif block["type"] == "code":
            lines.append("```" + block["language"])
            lines.append(block["text"])
            lines.append("```")
    return "\n".join(lines)


def markdown_to_pptx(markdown: str, out_path: Path) -> Path:
    """마크다운 전체 콘텐츠를 PowerPoint로 변환한다."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    blocks = _parse_markdown_blocks(markdown)
    prs = Presentation()

    if not blocks:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.5), PptxInches(9), PptxInches(1)
        ).text_frame.text = "변환할 콘텐츠가 없습니다."
        prs.save(out_path)
        return out_path

    # 제목을 기준으로 슬라이드 그룹화
    slides: list[dict] = []
    current_title = ""
    current_body: list[dict] = []

    for block in blocks:
        if block["type"] == "heading":
            if current_title or current_body:
                slides.append({"title": current_title, "body": current_body})
            current_title = block["text"]
            current_body = []
        else:
            current_body.append(block)

    if current_title or current_body:
        slides.append({"title": current_title, "body": current_body})

    # 제목이 하나도 없으면 모든 콘텐츠를 한 슬라이드에 담는다
    if not slides:
        slides = [{"title": "", "body": blocks}]

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 제목
        title_box = slide.shapes.add_textbox(
            PptxInches(0.5), PptxInches(0.3), PptxInches(9), PptxInches(0.8)
        )
        title_box.text_frame.text = slide_data["title"]
        for paragraph in title_box.text_frame.paragraphs:
            paragraph.font.size = Pt(20)
            paragraph.font.bold = True

        # 본문
        body_text = _blocks_to_slide_text(slide_data["body"])
        if body_text:
            body_box = slide.shapes.add_textbox(
                PptxInches(0.5), PptxInches(1.2), PptxInches(9), PptxInches(5.5)
            )
            tf = body_box.text_frame
            tf.text = body_text
            tf.word_wrap = True
            for paragraph in tf.paragraphs:
                paragraph.font.size = Pt(12)

    prs.save(out_path)
    return out_path


# ---------------------------------------------------------------------------
# Markdown table extraction (basic conversion)
# ---------------------------------------------------------------------------

def _parse_markdown_tables(markdown: str) -> list[dict]:
    """마크다운에서 표만 추출한다.

    반환 형식: {"headers": [...], "rows": [[...], ...], "has_header": bool}
    """
    # 기존 로직과 동일
    lines = markdown.splitlines()
    tables: list[dict] = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        if not stripped.startswith("|"):
            i += 1
            continue
        table_lines: list[str] = []
        while i < len(lines) and lines[i].strip().startswith("|"):
            table_lines.append(lines[i].strip())
            i += 1
        if len(table_lines) < 2:
            continue
        has_separator = any(
            re.match(r"^\|?[\s\-:|]+\|?$", line) for line in table_lines
        )
        data_lines = [
            line for line in table_lines if not re.match(r"^\|?[\s\-:|]+\|?$", line)
        ]
        if not data_lines:
            continue
        first_row = [c.strip() for c in data_lines[0].split("|")[1:-1]]
        rows = [
            [c.strip() for c in line.split("|")[1:-1]]
            for line in data_lines[1:]
        ]
        rows = [r for r in rows if r]
        if has_separator and _is_valid_header(first_row):
            headers = first_row
            has_header = True
        else:
            # 헤더가 없는 연속 페이지: 첫 행도 데이터로 취급
            headers = []
            rows = [first_row] + rows
            has_header = False
        if headers or rows:
            tables.append({"headers": headers, "rows": rows, "has_header": has_header})
    return tables


def _parse_html_tables(markdown: str) -> list[dict]:
    """마크다운 내에 포함된 HTML <table> 블록을 표로 추출한다.

    OCR 파이프라인이 마크다운 대신 HTML 테이블을 생성할 때 사용된다.
    반환 형식: {"headers": [...], "rows": [[...], ...], "has_header": bool}.
    """
    tables: list[dict] = []
    # <table ...> ... </table> 블록 추출 (base64 이미지 등 무시)
    for table_html in re.findall(r"<table[^>]*>.*?</table>", markdown, flags=re.DOTALL | re.IGNORECASE):
        rows: list[list[str]] = []
        for tr in re.findall(r"<tr[^>]*>.*?</tr>", table_html, flags=re.DOTALL | re.IGNORECASE):
            cells: list[str] = []
            for cell in re.findall(r"<t[dh][^>]*>(.*?)</t[dh]>", tr, flags=re.DOTALL | re.IGNORECASE):
                # 내부 HTML 태그 제거, HTML 엔티티 디코드, 공백 정리
                text = re.sub(r"<[^>]+>", "", cell)
                text = html_module.unescape(text)
                cells.append(text.strip())
            if cells:
                rows.append(cells)
        if not rows:
            continue
        # 첫 행이 실제 헤더인지 판단
        if _is_valid_header(rows[0]):
            tables.append({"headers": rows[0], "rows": rows[1:], "has_header": True})
        else:
            tables.append({"headers": [], "rows": rows, "has_header": False})
    return tables


def _extract_tables(markdown: str) -> list[dict]:
    """마크다운 표와 HTML 표를 모두 추출하여 하나의 목록으로 반환한다.

    [Flow: Step 1 (마크다운 표 파싱) -> Step 2 (HTML 표 파싱) -> Step 3 (순서대로 결합)]
    """
    return _parse_markdown_tables(markdown) + _parse_html_tables(markdown)


def _majority_col_count(rows: list[list[str]]) -> int:
    """행 목록에서 가장 많이 나타나는 열 수를 반환한다."""
    if not rows:
        return 0
    counts: dict[int, int] = {}
    for row in rows:
        counts[len(row)] = counts.get(len(row), 0) + 1
    return max(counts.items(), key=lambda item: item[1])[0]


def _resolve_merge_threshold(threshold: float | None) -> float:
    """병합 유사도 임계값을 결정한다. 인자가 없으면 설정값을 사용한다."""
    if threshold is not None:
        return threshold
    try:
        from ..config import settings
        return float(settings.table_merge_similarity_threshold)
    except Exception:
        # 설정 로드 실패 시 안전한 기본값
        return 0.75


def _table_similarity(prev: dict, table: dict) -> float:
    """이전 표와 현재 표의 형식 유사도를 0~1로 산출한다.

    [Flow: Step 1 (이전 표 열 프로파일 구축) -> Step 2 (현재 표 데이터행 정렬 confidence 평균)
            -> Step 3 (헤더 문자열 유사도) -> Step 4 (가중 결합 점수 반환)]

    데이터 점수(정렬 적합도)와 헤더 문자열 유사도를 결합한다.
    헤더 정보가 없으면 데이터 점수만 사용한다.
    """
    from difflib import SequenceMatcher

    prev_headers = prev.get("headers", [])
    prev_rows = prev.get("rows", [])
    cur_headers = table.get("headers", [])
    cur_rows = table.get("rows", [])
    prev_col_count = len(prev_headers)
    if prev_col_count == 0 or not cur_rows:
        return 0.0

    # 열 수 차이가 큰 표(±2 초과)는 다른 표로 간주하여 병합 후보 제외
    cur_col_count = _majority_col_count(cur_rows)
    if abs(cur_col_count - prev_col_count) > 2:
        return 0.0

    # Step 1~2: 이전 표 프로파일에 현재 표 데이터행을 정렬해 confidence 평균 계산
    profiles = _build_column_profiles(prev_headers, prev_rows)
    sample_rows = cur_rows[:20]  # 성능을 위해 상위 20행만 표본 사용
    confidences = [_align_row_to_columns(row, profiles)[1] for row in sample_rows]
    data_score = sum(confidences) / len(confidences) if confidences else 0.0

    # Step 3: 헤더 문자열 유사도 (양쪽 헤더가 있고 열 수가 같을 때만)
    header_score = None
    if cur_headers and prev_col_count == len(cur_headers):
        prev_text = " ".join(prev_headers)
        cur_text = " ".join(cur_headers)
        header_score = SequenceMatcher(None, prev_text, cur_text).ratio()

    # Step 4: 가중 결합 (헤더 점수가 없으면 데이터 점수만)
    if header_score is None:
        return data_score
    return 0.7 * data_score + 0.3 * header_score


def _merge_tables(tables: list[dict], threshold: float | None = None) -> list[dict]:
    """동일 헤더를 가진 표와 헤더 없는 연속 표를 통합한다.

    페이지 마커(<!-- 페이지 N -->) 등으로 분리된 표 중,
    헤더가 같거나 후속 페이지가 헤더 없이 데이터 행만 있는 경우,
    또는 형식 유사도가 임계값 이상인 연속 표를 하나의 표로 통합한다.
    유사도 병합 시 이전 표 헤더를 유지하고 현재 표 헤더는 데이터로 흡수한다.
    """
    if not tables:
        return []
    sim_threshold = _resolve_merge_threshold(threshold)
    merged: list[dict] = []
    for table in tables:
        has_header = table.get("has_header", True)
        headers = table.get("headers", [])
        rows = table.get("rows", [])
        if not merged:
            if has_header and headers:
                merged.append({"headers": list(headers), "rows": [list(r) for r in rows]})
            else:
                # 첫 표부터 헤더가 없으면 첫 데이터 행을 헤더로 승격 (fallback)
                all_rows = ([list(headers)] + [list(r) for r in rows]) if headers else [list(r) for r in rows]
                if all_rows:
                    merged.append({"headers": all_rows[0], "rows": all_rows[1:]})
                else:
                    merged.append({"headers": [], "rows": []})
            continue
        last = merged[-1]
        last_col_count = len(last["headers"])
        if has_header and headers:
            if tuple(headers) == tuple(last["headers"]):
                last["rows"].extend([list(r) for r in rows])
            elif _table_similarity(last, table) >= sim_threshold:
                # 형식 유사: 이전 표 헤더를 유지하고 현재 표 헤더는 버림(중복 헤더로 간주), 데이터 행만 병합
                last["rows"].extend([list(r) for r in rows])
            else:
                merged.append({"headers": list(headers), "rows": [list(r) for r in rows]})
        else:
            # 헤더 없는 연속 표: 이전 표의 헤더를 상속받아 병합 (가장 흔한 케이스)
            if rows:
                current_col_count = _majority_col_count(rows)
                # 열 수가 같거나, 열 수가 달라도 형식 유사도가 임계값 이상이면 병합
                if current_col_count == last_col_count or _table_similarity(last, table) >= sim_threshold:
                    last["rows"].extend([list(r) for r in rows])
                else:
                    # 형식이 다르면 독립 표로 시작 (첫 행을 헤더로 승격)
                    merged.append({"headers": list(rows[0]), "rows": [list(r) for r in rows[1:]]})
            else:
                merged.append({"headers": list(headers), "rows": []})
    return merged


def _normalize_rows(tables: list[dict]) -> list[dict]:
    """각 표의 모든 행을 동일한 컬럼 수로 맞춘다."""
    for table in tables:
        col_count = max(len(table["headers"]), max((len(row) for row in table["rows"]), default=0))
        table["headers"] = table["headers"] + [""] * (col_count - len(table["headers"]))
        for row in table["rows"]:
            row.extend([""] * (col_count - len(row)))
    return tables


# ---------------------------------------------------------------------------
# 데이터 형식 기반 열 정렬 보정 (basic conversion)
# [Flow: Step 1 (셀 형식 분류) -> Step 2 (정상 행으로 열 프로파일 구축)
#        -> Step 3 (열 수 어긋난 행을 순차 정렬로 보정) -> Step 4 (보정 로그 수집)]
# ---------------------------------------------------------------------------

# 형식 판별용 정규식 (구체적 형식 우선)
_RE_EMAIL = re.compile(r"^[\w.+-]+@[\w-]+\.[\w.-]+$")
_RE_RRN = re.compile(r"^\d{6}\s*-\s*\d{7}$")  # 주민등록번호
_RE_BIZNO = re.compile(r"^\d{3}\s*-\s*\d{2}\s*-\s*\d{5}$")  # 사업자등록번호
_RE_PHONE = re.compile(r"^(0\d{1,2}|\+?\d{1,3})[\s-]?\d{3,4}[\s-]?\d{4}$")
_RE_ACCOUNT = re.compile(r"^\d{2,6}([\s-]\d{2,6}){1,4}$")  # 계좌번호 (구분자 포함 숫자 그룹)
_RE_DATE = re.compile(
    r"^\d{4}[.\-/]\s?\d{1,2}[.\-/]\s?\d{1,2}\.?$"
    r"|^\d{4}년\s?\d{1,2}월\s?\d{1,2}일$"
    r"|^\d{1,2}[.\-/]\d{1,2}$"
)
_RE_AMOUNT = re.compile(r"^[₩$]?\s?-?\d{1,3}(,\d{3})+(\.\d+)?\s?(원|KRW|USD)?$")
_RE_INTEGER = re.compile(r"^-?\d+$")
_RE_DECIMAL = re.compile(r"^-?\d+\.\d+$")
_RE_KOREAN_NAME = re.compile(r"^[가-힣]{2,4}$")


def _classify_cell(value: str) -> str:
    """셀 문자열을 데이터 형식으로 분류한다.

    반환값: empty, email, rrn, bizno, phone, account, date, amount,
            decimal, integer, korean_name, string 중 하나.
    구체적인 형식을 먼저 검사하고, 일치하지 않으면 일반 문자열로 처리한다.
    """
    text = (value or "").strip()
    if not text:
        return "empty"
    if _RE_EMAIL.match(text):
        return "email"
    if _RE_RRN.match(text):
        return "rrn"
    if _RE_BIZNO.match(text):
        return "bizno"
    if _RE_PHONE.match(text):
        return "phone"
    if _RE_DATE.match(text):
        return "date"
    if _RE_AMOUNT.match(text):
        return "amount"
    if _RE_DECIMAL.match(text):
        return "decimal"
    if _RE_INTEGER.match(text):
        # 계좌번호 형태(구분자 포함)는 위에서 account로 처리되지 않으므로 여기서 순수 정수만
        return "integer"
    if _RE_ACCOUNT.match(text):
        return "account"
    if _RE_KOREAN_NAME.match(text):
        return "korean_name"
    return "string"


def _is_valid_header(row: list[str]) -> bool:
    """표 헤더 행이 실제로 유효한 헤더인지 판단한다.

    데이터 행(금액, 날짜, 순번 등)이 헤더로 잘못 인식되는 것을 막는다.
    """
    if not row:
        return False
    non_empty = [cell for cell in row if cell.strip()]
    if not non_empty:
        return False
    data_like_types = {"integer", "decimal", "amount", "date", "rrn", "bizno", "phone", "account"}
    data_like_count = sum(1 for cell in non_empty if _classify_cell(cell) in data_like_types)
    return data_like_count <= len(non_empty) / 2


# 형식 그룹: 같은 그룹 내 형식은 부분 일치로 간주하여 페널티를 낮춘다.
_TYPE_GROUP = {
    "integer": "numeric",
    "decimal": "numeric",
    "amount": "numeric",
    "account": "numeric",
    "phone": "numeric",
    "rrn": "numeric",
    "bizno": "numeric",
    "korean_name": "text",
    "string": "text",
    "email": "text",
    "date": "text",
    "empty": "empty",
}


def _build_column_profiles(headers: list[str], rows: list[list[str]]) -> list[dict]:
    """열 수가 헤더와 정확히 일치하는 '정상 행'만으로 각 열의 형식 프로파일을 만든다.

    반환: 각 열별 {"dominant": 형식, "confidence": 비율, "counts": {형식: 개수}}.
    정상 행이 3개 미만이면 신뢰도가 낮으므로 confidence를 0으로 표시한다.
    """
    from collections import Counter

    col_count = len(headers)
    clean_rows = [r for r in rows if len(r) == col_count]
    profiles: list[dict] = []
    for col_idx in range(col_count):
        counter: Counter = Counter()
        for row in clean_rows:
            # 빈칸은 형식 판단에 약하게만 반영 (모든 열이 가끔 빌 수 있음)
            counter[_classify_cell(row[col_idx])] += 1
        non_empty_total = sum(c for t, c in counter.items() if t != "empty")
        if counter:
            dominant, dom_count = counter.most_common(1)[0]
            # 지배 형식이 empty면 non-empty 중 최빈값을 우선 사용
            if dominant == "empty" and non_empty_total > 0:
                non_empty = [(t, c) for t, c in counter.items() if t != "empty"]
                dominant, dom_count = max(non_empty, key=lambda x: x[1])
            total = sum(counter.values())
            confidence = dom_count / total if total else 0.0
        else:
            dominant, confidence = "string", 0.0
        if len(clean_rows) < 3:
            confidence = 0.0
        profiles.append({"dominant": dominant, "confidence": confidence, "counts": dict(counter)})
    return profiles


def _type_match_cost(cell_type: str, expected_type: str) -> float:
    """셀 형식과 기대 형식 간의 불일치 비용(0=완전일치, 1=완전불일치)을 계산한다."""
    if cell_type == expected_type:
        return 0.0
    if cell_type == "empty" or expected_type == "empty":
        return 0.3  # 빈칸은 어느 열에나 어느 정도 허용
    if _TYPE_GROUP.get(cell_type) == _TYPE_GROUP.get(expected_type):
        return 0.4  # 같은 그룹(숫자류/문자류)이면 부분 일치
    return 1.0


_GAP_PENALTY = 0.6   # 열을 빈칸으로 건너뛰는 비용 (M<N)
_MERGE_PENALTY = 0.5  # 인접 셀을 병합하는 비용 (M>N)


def _align_row_to_columns(cells: list[str], profiles: list[dict]) -> tuple[list[str], float]:
    """관측 셀 M개를 기준 열 N개에 순서를 보존하며 정렬한다 (시퀀스 정렬 DP).

    연산: 매칭(형식 비용), 빈칸 삽입(M<N, gap 페널티),
          인접 셀 병합(M>N, merge 페널티 후 병합 셀 형식 재평가).
    반환: (정렬된 N길이 행, confidence 0~1).
    """
    m = len(cells)
    n = len(profiles)
    if n == 0:
        return [], 0.0
    if m == n:
        # 열 수 일치: 정렬 불필요, 형식 적합도만 계산
        conf = _row_confidence(cells, profiles)
        return list(cells), conf

    INF = float("inf")
    # dp[i][j] = 관측 i개, 기준열 j개까지 정렬한 최소 비용
    dp = [[INF] * (n + 1) for _ in range(m + 1)]
    # back[i][j] = (op, prev_i, prev_j, 배치된 값 or None)
    back: list[list] = [[None] * (n + 1) for _ in range(m + 1)]
    dp[0][0] = 0.0

    for i in range(m + 1):
        for j in range(n + 1):
            if dp[i][j] == INF:
                continue
            cur = dp[i][j]
            # 1) 빈칸 삽입: 기준열 j를 빈칸으로 채우고 다음 열로 (관측 소비 없음)
            if j < n:
                cost = cur + _GAP_PENALTY
                if cost < dp[i][j + 1]:
                    dp[i][j + 1] = cost
                    back[i][j + 1] = ("gap", i, j, "")
            # 2) 매칭: 관측 i를 기준열 j에 배치
            if i < m and j < n:
                cell_type = _classify_cell(cells[i])
                cost = cur + _type_match_cost(cell_type, profiles[j]["dominant"])
                if cost < dp[i + 1][j + 1]:
                    dp[i + 1][j + 1] = cost
                    back[i + 1][j + 1] = ("match", i, j, cells[i])
            # 3) 병합: 관측 i,i+1을 합쳐 기준열 j에 배치 (M>N 처리)
            if i + 1 < m and j < n:
                merged = (cells[i] + " " + cells[i + 1]).strip()
                merged_type = _classify_cell(merged)
                cost = cur + _MERGE_PENALTY + _type_match_cost(merged_type, profiles[j]["dominant"])
                if cost < dp[i + 2][j + 1]:
                    dp[i + 2][j + 1] = cost
                    back[i + 2][j + 1] = ("merge", i, j, merged)

    if dp[m][n] == INF:
        # 정렬 실패: 잘라내거나 빈칸 패딩으로 fallback
        aligned = (list(cells) + [""] * n)[:n]
        return aligned, 0.0

    # 역추적으로 정렬된 행 복원
    aligned = [""] * n
    i, j = m, n
    while back[i][j] is not None:
        op, pi, pj, val = back[i][j]
        if op in ("match", "merge"):
            aligned[pj] = val
        i, j = pi, pj

    conf = _row_confidence(aligned, profiles)
    return aligned, conf


def _row_confidence(row: list[str], profiles: list[dict]) -> float:
    """정렬된 행이 열 프로파일과 얼마나 잘 맞는지 0~1로 산출한다."""
    if not profiles:
        return 0.0
    total = 0.0
    for idx, prof in enumerate(profiles):
        cell = row[idx] if idx < len(row) else ""
        cost = _type_match_cost(_classify_cell(cell), prof["dominant"])
        # 프로파일 신뢰도로 가중 (신뢰도 낮은 열은 판단 근거 약함)
        weight = 0.5 + 0.5 * prof.get("confidence", 0.0)
        total += (1.0 - cost) * weight
    max_score = sum(0.5 + 0.5 * p.get("confidence", 0.0) for p in profiles)
    return total / max_score if max_score else 0.0


def _correct_tables(tables: list[dict]) -> tuple[list[dict], list[dict]]:
    """열 수가 어긋난 행을 형식 기반으로 보정하고 보정 로그를 반환한다.

    반환: (보정된 tables, correction_log).
    correction_log 각 항목: {table_idx, row_num, original, corrected, confidence, reason}.
    """
    correction_log: list[dict] = []
    for table_idx, table in enumerate(tables):
        headers = table["headers"]
        col_count = len(headers)
        if col_count == 0:
            continue
        profiles = _build_column_profiles(headers, table["rows"])
        new_rows: list[list[str]] = []
        for row_num, row in enumerate(table["rows"], start=1):
            if len(row) == col_count:
                new_rows.append(row)
                continue
            aligned, confidence = _align_row_to_columns(row, profiles)
            new_rows.append(aligned)
            correction_log.append({
                "table_idx": table_idx + 1,
                "row_num": row_num,
                "original": " | ".join(row),
                "corrected": " | ".join(aligned),
                "confidence": round(confidence, 2),
                "reason": f"열 수 불일치 ({len(row)}→{col_count})",
            })
        table["rows"] = new_rows
    return tables, correction_log


def markdown_to_csv_basic(markdown: str, out_path: Path) -> Path:
    """마크다운 표를 하나의 CSV 파일로 통합한다.

    형식 기반 열 정렬 보정을 적용한다. (보정 로그는 xlsx에만 기록하고 CSV는 보정만 적용)
    """
    import csv as csv_module

    out_path.parent.mkdir(parents=True, exist_ok=True)
    corrected, _log = _correct_tables(_merge_tables(_extract_tables(markdown)))
    tables = _normalize_rows(corrected)
    with open(out_path, "w", newline="", encoding="utf-8-sig") as f:
        writer = csv_module.writer(f)
        for idx, table in enumerate(tables):
            if idx > 0:
                writer.writerow([])
            writer.writerow(table["headers"])
            for row in table["rows"]:
                writer.writerow(row)
    return out_path


def markdown_to_xlsx_basic(markdown: str, out_path: Path) -> Path:
    """마크다운 표를 하나의 xlsx 파일에 통합한다.

    형식 기반 열 정렬 보정을 적용하고, 보정된 행이 있으면 '보정 내역' 시트를 추가한다.
    헤더가 다른 표는 별도 워크시트('table1', 'table2'...)로 분리된다.
    """
    out_path.parent.mkdir(parents=True, exist_ok=True)
    corrected, correction_log = _correct_tables(_merge_tables(_extract_tables(markdown)))
    tables = _normalize_rows(corrected)
    wb = Workbook()
    if tables:
        default_sheet = wb.active
        if default_sheet is not None:
            wb.remove(default_sheet)
        for idx, table in enumerate(tables, start=1):
            sheet_name = _safe_sheet_name(f"table{idx}", idx)
            ws = wb.create_sheet(title=sheet_name)
            _write_table_to_sheet(ws, table["headers"], table["rows"])
    else:
        default_sheet = wb.active
        if default_sheet is not None:
            default_sheet.title = "Tables"

    # 보정된 행이 있으면 '보정 내역' 시트를 추가하여 사용자가 검토할 수 있게 한다
    if correction_log:
        _write_correction_sheet(wb, correction_log)

    wb.save(out_path)
    return out_path


def _write_table_to_sheet(ws, headers: list[str], rows: list[list[str]]) -> None:
    """표의 헤더와 데이터 행을 워크시트에 작성하고 컬럼 너비를 조정한다."""
    for col_idx, h in enumerate(headers):
        cell = ws.cell(1, col_idx + 1, h)
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="E0E7FF", end_color="E0E7FF", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    for row_idx, table_row in enumerate(rows, start=2):
        for col_idx, val in enumerate(table_row):
            ws.cell(row_idx, col_idx + 1, val)

    # 컬럼 너비 자동 조정
    for col in ws.columns:
        max_len = 0
        col_letter = col[0].column_letter
        for cell in col:
            try:
                if cell.value:
                    max_len = max(max_len, len(str(cell.value)))
            except Exception:
                pass
        ws.column_dimensions[col_letter].width = min(max_len + 2, 60)


def _write_correction_sheet(wb, correction_log: list[dict]) -> None:
    """보정 로그를 별도 시트에 기록한다. (열 수 불일치로 형식 기반 재정렬된 행 목록)"""
    ws = wb.create_sheet("보정 내역")
    log_headers = ["표 번호", "행 번호", "원본", "보정 결과", "신뢰도", "사유"]
    for col_idx, h in enumerate(log_headers):
        cell = ws.cell(1, col_idx + 1, h)
        cell.font = Font(bold=True)
        cell.fill = PatternFill(start_color="FEE2E2", end_color="FEE2E2", fill_type="solid")
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
    for row_idx, entry in enumerate(correction_log, start=2):
        ws.cell(row_idx, 1, entry["table_idx"])
        ws.cell(row_idx, 2, entry["row_num"])
        ws.cell(row_idx, 3, entry["original"])
        ws.cell(row_idx, 4, entry["corrected"])
        ws.cell(row_idx, 5, entry["confidence"])
        ws.cell(row_idx, 6, entry["reason"])
    widths = [8, 8, 50, 50, 8, 20]
    for col_idx, w in enumerate(widths):
        ws.column_dimensions[get_column_letter(col_idx + 1)].width = w


# ---------------------------------------------------------------------------
# XLSX (legacy full markdown conversion)
# ---------------------------------------------------------------------------

def _safe_sheet_name(name: str, index: int) -> str:
    """openpyxl 시트 이름 규칙에 맞게 제한 길이 및 금지 문자를 처리한다."""
    name = re.sub(r"[\\/*?:\[\]]", "-", name).strip()[:28]
    if not name:
        name = f"Sheet{index}"
    return name


def markdown_to_xlsx(markdown: str, out_path: Path) -> Path:
    """마크다운 전체 콘텐츠를 Excel 파일로 변환한다."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    blocks = _parse_markdown_blocks(markdown)
    wb = Workbook()
    ws_content = wb.active
    ws_content.title = "Content"

    # 콘텐츠 시트에 텍스트/목록/제목/코드 기록
    row = 1
    table_count = 0
    for block in blocks:
        if block["type"] == "heading":
            cell = ws_content.cell(row, 1, block["text"])
            cell.font = Font(bold=True, size=16 - block["level"])
            row += 1
        elif block["type"] == "paragraph":
            ws_content.cell(row, 1, block["text"])
            row += 1
        elif block["type"] == "list":
            for idx, item in enumerate(block["items"], start=1):
                prefix = f"{idx}. " if block["ordered"] else "• "
                ws_content.cell(row, 1, prefix + item)
                row += 1
        elif block["type"] == "code":
            cell = ws_content.cell(row, 1, block["text"])
            cell.font = Font(name="Courier New")
            row += 1
        elif block["type"] == "table":
            table_count += 1
            headers = block["headers"]
            rows = block["rows"]
            sheet_name = _safe_sheet_name("Table " + " ".join(headers[:2]), table_count)
            ws_table = wb.create_sheet(title=sheet_name)
            for c_idx, h in enumerate(headers):
                cell = ws_table.cell(1, c_idx + 1, h)
                cell.font = Font(bold=True)
                cell.fill = PatternFill(start_color="E0E7FF", end_color="E0E7FF", fill_type="solid")
                cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            for r_idx, row_data in enumerate(rows):
                for c_idx, val in enumerate(row_data):
                    ws_table.cell(r_idx + 2, c_idx + 1, val)
            # 콘텐츠 시트에 표 위치 안내
            ws_content.cell(row, 1, f"[표 {table_count}: {sheet_name} 시트 참조]")
            row += 1

    # 표가 하나도 없으면 안내 문구 추가
    if table_count == 0 and row == 1:
        ws_content.cell(1, 1, "변환할 콘텐츠가 없습니다.")

    # 컬럼 너비 자동 조정
    for ws in wb.worksheets:
        for col in ws.columns:
            max_len = 0
            col_letter = col[0].column_letter
            for cell in col:
                try:
                    if cell.value:
                        max_len = max(max_len, len(str(cell.value)))
                except Exception:
                    pass
            ws.column_dimensions[col_letter].width = min(max_len + 2, 60)

    wb.save(out_path)
    return out_path
